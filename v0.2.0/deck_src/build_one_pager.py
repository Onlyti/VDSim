"""Build the VDSim 1-page deck (PNG + PPTX + PDF).

Layout (1 page, 16:9 landscape):
    [ Title bar: navy #002060 ]
    [ 1-line elevator pitch ]
    [ Left ]                          [ Right ]
     Ld × Lc capability matrix         Competitive radar (Adams, VI,
       (Ld1-Ld5 × Lc1-Lc8 grid,         CarMaker, CarSim, Simulink,
        24+ verified cells)             VDSim PoC, VDSim+L4-L5)
    [ Bottom band ]
     - Status (141 tests, 95% W1-W12, 4 vehicles, 8 scenarios)
     - Phase 2 roadmap (CARLA UE5, CarMaker ERG, Ld4-Ld5, MPC)

Outputs:
    docs/deck_src/poc_one_pager.png
    docs/deck_src/poc_one_pager.pptx
    docs/deck_src/poc_one_pager.pdf
"""
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import FancyBboxPatch

OUT = Path(__file__).resolve().parent
OUT.mkdir(parents=True, exist_ok=True)

# AILab palette
NAVY     = "#002060"
PRIMARY  = "#005195"
ACCENT   = "#01A0E9"
WARNING  = "#DC291E"
BODY     = "#333333"
GREY     = "#7c8693"
LIGHT_GREY = "#e6ebf0"

# ============== matplotlib one-pager ==============
fig = plt.figure(figsize=(16, 9), facecolor="white")

# Title bar
title_ax = fig.add_axes([0, 0.92, 1, 0.08])
title_ax.add_patch(patches.Rectangle((0, 0), 1, 1, color=NAVY))
title_ax.text(0.012, 0.50, "VDSim", color="white", fontsize=22, weight="bold",
              va="center", ha="left", family="DejaVu Sans")
title_ax.text(0.110, 0.50,
              "Open-Core Vehicle Dynamics Simulator  —  unified Lc × Ld ladder ABI",
              color="white", fontsize=15, va="center", ha="left",
              family="DejaVu Sans")
title_ax.text(0.985, 0.50, "PoC W1-W12 · 95 %", color=ACCENT, fontsize=12,
              va="center", ha="right", family="DejaVu Sans", weight="bold")
title_ax.axis("off")

# Pitch row
pitch_ax = fig.add_axes([0, 0.86, 1, 0.06])
pitch_ax.text(0.5, 0.5,
              "Bridges autonomy stack (throttle/pedal/path) and chassis design "
              "(hardpoint/multibody) in a single C++17 + Python ABI.",
              color=BODY, fontsize=12, va="center", ha="center",
              family="DejaVu Sans", style="italic")
pitch_ax.axis("off")

# ====== Left: Ld × Lc matrix ======
left_ax = fig.add_axes([0.03, 0.20, 0.46, 0.65])
left_ax.set_title("Capability matrix:  Ld dynamics × Lc control",
                  fontsize=13, color=NAVY, weight="bold", pad=10, loc="left")

ld_rows = ["Ld1-Bicycle (5 DOF)",
           "Ld2-SevenDOF (7)",
           "Ld3-FourteenDOF (14)",
           "Ld4-MultibodyKin", "Ld5-MultibodyCompliant"]
lc_cols = ["Lc1-PerWheel", "Lc2-Axle", "Lc3-FxTotal",
           "Lc4-Pedal", "Lc5-AxTgt", "Lc6-VTgt",
           "Lc7-PurePursuit", "Lc8-Waypoint"]

# 0=NA, 1=via lowering, 2=full primary, 3=planned
matrix = np.array([
    [1, 1, 1, 2, 1, 1, 1, 1],   # Ld1
    [1, 1, 1, 2, 1, 1, 1, 1],   # Ld2
    [1, 1, 1, 2, 1, 1, 1, 1],   # Ld3
    [3, 3, 3, 3, 3, 3, 3, 3],   # Ld4
    [3, 3, 3, 3, 3, 3, 3, 3],   # Ld5
])
color_for = {0: LIGHT_GREY, 1: ACCENT, 2: PRIMARY, 3: GREY}
label_for = {0: "—",        1: "via",   2: "✓",     3: "plan"}

cell_w = 1.0 / len(lc_cols)
cell_h = 1.0 / len(ld_rows)
for i in range(len(ld_rows)):
    for j in range(len(lc_cols)):
        v = matrix[i, j]
        rect = patches.Rectangle((j * cell_w, 1 - (i + 1) * cell_h),
                                 cell_w * 0.96, cell_h * 0.85,
                                 facecolor=color_for[v], edgecolor="white")
        left_ax.add_patch(rect)
        txt_color = "white" if v in (1, 2) else BODY
        left_ax.text(j * cell_w + cell_w / 2,
                     1 - (i + 1) * cell_h + cell_h * 0.45,
                     label_for[v], ha="center", va="center",
                     color=txt_color, fontsize=10, weight="bold")

# Row labels
for i, name in enumerate(ld_rows):
    color = NAVY if i < 3 else GREY
    left_ax.text(-0.01, 1 - (i + 1) * cell_h + cell_h * 0.45,
                 name, ha="right", va="center", color=color,
                 fontsize=10, weight="bold" if i < 3 else "normal")

# Col labels
for j, name in enumerate(lc_cols):
    left_ax.text(j * cell_w + cell_w / 2, 1.02, name,
                 ha="center", va="bottom", color=NAVY, fontsize=9,
                 weight="bold", rotation=25)

# Legend
legend_y = -0.16
for k, (v, lbl) in enumerate([(2, "✓ verified primary"),
                              (1, "via lowering"),
                              (3, "planned (Phase 2)")]):
    x = k * 0.34
    left_ax.add_patch(patches.Rectangle((x, legend_y - 0.02), 0.025, 0.04,
                                         color=color_for[v]))
    left_ax.text(x + 0.035, legend_y, lbl, color=BODY, fontsize=9,
                  va="center")
left_ax.text(0.5, legend_y - 0.06,
             "24 / 40 cells verified · 141 / 141 tests passing",
             color=PRIMARY, fontsize=10, weight="bold", ha="center")

left_ax.set_xlim(-0.30, 1.02)
left_ax.set_ylim(-0.20, 1.08)
left_ax.axis("off")

# ====== Right: Competitive comparison ======
right_ax = fig.add_axes([0.55, 0.32, 0.42, 0.55], polar=True)
axes_labels = [
    "Multibody\ndepth",
    "Simplified VD\n(autonomy)",
    "Control ladder\n(L1-L8)",
    "Open-core\nAPI",
    "External\nintegration",
    "Tire model\nfidelity",
    "Free /\nOpen license",
]
N = len(axes_labels)
angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
angles += angles[:1]

solutions = {
    "Adams Car":              ([5, 1, 1, 0, 2, 5, 0], "#7F0000", 1.0),
    "VI-CarRealTime":         ([5, 2, 3, 0, 3, 4, 0], "#DC291E", 1.0),
    "CarMaker":               ([2, 5, 3, 0, 3, 5, 1], "#345A8A", 1.0),
    "CarSim":                 ([2, 5, 3, 0, 3, 4, 1], "#4F81BD", 1.0),
    "Simulink VDB":           ([1, 4, 3, 1, 5, 3, 2], "#01A0E9", 1.0),
    "VDSim (current PoC)":    ([1, 5, 5, 5, 4, 3, 5], "#FF8C00", 2.5),
    "VDSim (+ L4-L5 plan)":   ([4, 5, 5, 5, 5, 4, 5], "#FFD700", 2.5),
}
for name, (vals, color, lw) in solutions.items():
    vals = vals + vals[:1]
    right_ax.plot(angles, vals, "o-", linewidth=lw, color=color,
                   label=name, markersize=3)
    if "VDSim" in name:
        right_ax.fill(angles, vals, alpha=0.10, color=color)
right_ax.set_xticks(angles[:-1])
right_ax.set_xticklabels(axes_labels, fontsize=8)
right_ax.set_ylim(0, 5)
right_ax.set_yticks([1, 2, 3, 4, 5])
right_ax.set_yticklabels(["1", "2", "3", "4", "5"], fontsize=7)
right_ax.set_title("Competitive matrix vs commercial",
                   fontsize=13, color=NAVY, weight="bold", pad=15)
right_ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.30),
                fontsize=7, ncol=2, frameon=False)

# ====== Bottom: status + roadmap ======
bottom_ax = fig.add_axes([0, 0.04, 1, 0.15])
bottom_ax.add_patch(patches.Rectangle((0, 0), 1, 1, color=LIGHT_GREY))
bottom_ax.text(0.01, 0.78, "Status (now)",
               fontsize=11, weight="bold", color=NAVY, va="center")
bottom_ax.text(0.01, 0.46,
               "· 141 / 141 tests passing · 4 vehicles (sedan, sports, FSK, race) · 8 scenarios · 10 demo binaries",
               fontsize=9.5, color=BODY, va="center")
bottom_ax.text(0.01, 0.22,
               "· Lc4-Pedal → Lc8-Waypoint full cascade · Ld1-Ld3 full · CARLA-ready raycast ABI · pybind11 + 3D viewer",
               fontsize=9.5, color=BODY, va="center")

bottom_ax.axvline(0.50, ymin=0.15, ymax=0.85, color=NAVY, lw=1)

bottom_ax.text(0.52, 0.78, "Next (Phase 2)",
               fontsize=11, weight="bold", color=NAVY, va="center")
bottom_ax.text(0.52, 0.46,
               "· CARLA UE5 host integration · CarMaker ERG cross-validation · Ld4 MacPherson FK",
               fontsize=9.5, color=BODY, va="center")
bottom_ax.text(0.52, 0.22,
               "· SMPC / MPC (HPIPM) for path tracking · MF2002 tire + camber-Mz · vehicle calibration (TUR / FSK)",
               fontsize=9.5, color=BODY, va="center")

bottom_ax.set_xlim(0, 1); bottom_ax.set_ylim(0, 1); bottom_ax.axis("off")

# Footer
foot_ax = fig.add_axes([0, 0, 1, 0.035])
foot_ax.add_patch(patches.Rectangle((0, 0), 1, 1, color=NAVY))
foot_ax.text(0.012, 0.5, "Hanyang AILab · VDSim PoC summary v3",
             color="white", fontsize=8.5, va="center", ha="left")
foot_ax.text(0.988, 0.5,
             "github.com/Onlyti/VDSim · branch feat/poc-w5-to-95pct",
             color="white", fontsize=8.5, va="center", ha="right")
foot_ax.axis("off")

plt.savefig(OUT / "poc_one_pager.png", dpi=160, bbox_inches=None,
            facecolor="white")
plt.savefig(OUT / "poc_one_pager.pdf", bbox_inches=None, facecolor="white")
plt.close(fig)
print("PNG + PDF written")

# ============== python-pptx version ==============
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def navy_rgb():     return RGBColor(0x00, 0x20, 0x60)
def primary_rgb():  return RGBColor(0x00, 0x51, 0x95)
def accent_rgb():   return RGBColor(0x01, 0xA0, 0xE9)
def body_rgb():     return RGBColor(0x33, 0x33, 0x33)
def light_rgb():    return RGBColor(0xE6, 0xEB, 0xF0)

# Title bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0.05),
                              Inches(13.33), Inches(0.6))
bar.fill.solid(); bar.fill.fore_color.rgb = navy_rgb()
bar.line.fill.background()
tb = bar.text_frame; tb.margin_left = Inches(0.20); tb.margin_top = Inches(0.05)
p = tb.paragraphs[0]
p.text = "VDSim    Open-Core Vehicle Dynamics Simulator"
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.size = Pt(22); p.font.bold = True; p.font.name = "Tahoma"

# Pitch
tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.70),
                               Inches(12.3), Inches(0.4))
tf = tx.text_frame; p = tf.paragraphs[0]
p.text = ("Bridges autonomy stack (throttle / pedal / path) and chassis design "
          "(hardpoint / multibody) in a single C++17 + Python ABI.")
p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = body_rgb()
p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# Left + right images (matplotlib already rendered to PNG; embed)
# Slide approach: put the whole PNG as the slide. Editable text on top.
slide.shapes.add_picture(str(OUT / "poc_one_pager.png"),
                          Inches(0), Inches(0), width=Inches(13.33),
                          height=Inches(7.5))

# Footer overlay note (so the pptx can be edited; the PNG is the background)
note = slide.notes_slide.notes_text_frame
note.text = ("VDSim PoC W1-W12 95% summary. "
              "Editable PNG-backed slide; replace with native shapes "
              "before printing to PDF for editable layers.")

prs.save(OUT / "poc_one_pager.pptx")
print("PPTX written")

# Convert PPTX → PDF (LibreOffice)
import subprocess
res = subprocess.run(
    ["libreoffice", "--headless", "--convert-to", "pdf",
     "--outdir", str(OUT), str(OUT / "poc_one_pager.pptx")],
    capture_output=True, text=True, timeout=60)
print("libreoffice pdf:", res.returncode, res.stdout.splitlines()[-1] if res.stdout else "")

print("done")
